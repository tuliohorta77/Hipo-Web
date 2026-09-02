"""
HIPO — Testes do router /crm de propostas.

As regras puras estão em test_proposta_regras.py. Aqui o foco é o que só
aparece com banco: a numeração de versões, o snapshot de executivo e
cliente, o efeito colateral na mensalidade da oportunidade e o download.

O download monta o .pptx de verdade a partir de api/templates — é o único
teste que prova que o modelo versionado ainda tem os marcadores que o
código espera. Modelo trocado sem marcador gera proposta com campo em
branco, e ninguém percebe até o cliente receber.
"""
import uuid
from decimal import Decimal

import pytest

from services import proposta as regras
from services import proposta_render as render
from tests.conftest import criar_usuario

CNPJ_A = "11.222.333/0001-81"


# ── Helpers ──────────────────────────────────────────────────────────

async def nova_conta(client, headers, cnpj=CNPJ_A, razao="Metalurgica Alfa LTDA"):
    resp = await client.post(
        "/crm/contas", json={"razao_social": razao, "cnpj": cnpj}, headers=headers
    )
    assert resp.status_code == 201, resp.text
    return resp.json()


async def nova_oportunidade(client, headers, conta_id, **extra):
    corpo = {"conta_id": conta_id}
    corpo.update(extra)
    resp = await client.post("/crm/oportunidades", json=corpo, headers=headers)
    assert resp.status_code == 201, resp.text
    return resp.json()


def corpo_proposta(**troca):
    base = {
        "vidas": 50,
        "valor_por_vida": "20.00",
        "treinamentos": "2000.00",
        "laudos": "1000.00",
        "escopo": ["PGR - (NR-01)", "PCMSO - (NR-07)"],
        "data_proposta": "2026-08-26",
        "validade": "2026-09-05",
        "cidade": "Guarulhos",
    }
    base.update(troca)
    return base


async def cenario(client, headers):
    conta = await nova_conta(client, headers)
    opp = await nova_oportunidade(client, headers, conta["id"])
    return conta, opp


# ── Criação e versionamento ──────────────────────────────────────────

class TestCriar:
    async def test_primeira_e_versao_1(self, db_conn, client, usuario_adm):
        _, opp = await cenario(client, usuario_adm["headers"])
        resp = await client.post(
            f"/crm/oportunidades/{opp['id']}/propostas",
            json=corpo_proposta(), headers=usuario_adm["headers"],
        )
        assert resp.status_code == 201, resp.text
        assert resp.json()["versao"] == 1

    async def test_derivados_vem_calculados(self, db_conn, client, usuario_adm):
        """
        Mensalidade e investimento não são colunas: se a API não os
        calculasse na resposta, a tela mostraria o quadro vazio logo depois
        de gerar.
        """
        _, opp = await cenario(client, usuario_adm["headers"])
        body = (await client.post(
            f"/crm/oportunidades/{opp['id']}/propostas",
            json=corpo_proposta(), headers=usuario_adm["headers"],
        )).json()
        assert float(body["mensalidade"]) == 1000.0
        assert float(body["investimento"]) == 4000.0

    async def test_segunda_geracao_vira_v2(self, db_conn, client, usuario_adm):
        _, opp = await cenario(client, usuario_adm["headers"])
        for esperado in (1, 2, 3):
            body = (await client.post(
                f"/crm/oportunidades/{opp['id']}/propostas",
                json=corpo_proposta(), headers=usuario_adm["headers"],
            )).json()
            assert body["versao"] == esperado

    async def test_versao_e_por_oportunidade(self, db_conn, client, usuario_adm):
        """Cada negociação tem a própria v1 — a numeração não é global."""
        h = usuario_adm["headers"]
        conta = await nova_conta(client, h)
        a = await nova_oportunidade(client, h, conta["id"])
        b = await nova_oportunidade(client, h, conta["id"])
        await client.post(f"/crm/oportunidades/{a['id']}/propostas",
                          json=corpo_proposta(), headers=h)
        body = (await client.post(f"/crm/oportunidades/{b['id']}/propostas",
                                  json=corpo_proposta(), headers=h)).json()
        assert body["versao"] == 1

    async def test_atualiza_a_mensalidade_da_oportunidade(
        self, db_conn, client, usuario_adm
    ):
        """
        O funil soma ticket. Proposta enviada por R$ 1.000 com o funil
        marcando outro valor faz a previsão do mês mentir.
        """
        _, opp = await cenario(client, usuario_adm["headers"])
        await client.post(f"/crm/oportunidades/{opp['id']}/propostas",
                          json=corpo_proposta(), headers=usuario_adm["headers"])
        atual = (await client.get(f"/crm/oportunidades/{opp['id']}",
                                  headers=usuario_adm["headers"])).json()
        assert float(atual["valor_mensalidade"]) == 1000.0

    async def test_oportunidade_inexistente(self, db_conn, client, usuario_adm):
        resp = await client.post(
            f"/crm/oportunidades/{uuid.uuid4()}/propostas",
            json=corpo_proposta(), headers=usuario_adm["headers"],
        )
        assert resp.status_code == 404


# ── Snapshot ─────────────────────────────────────────────────────────

class TestSnapshot:
    async def test_copia_cliente_e_executivo(self, db_conn, client, usuario_adm):
        _, opp = await cenario(client, usuario_adm["headers"])
        body = (await client.post(
            f"/crm/oportunidades/{opp['id']}/propostas",
            json=corpo_proposta(), headers=usuario_adm["headers"],
        )).json()
        assert body["cliente_razao_social"] == "Metalurgica Alfa LTDA"
        assert body["executivo_nome"] == "Test ADM"

    async def test_renomear_a_conta_nao_reescreve_a_proposta(
        self, db_conn, client, usuario_adm
    ):
        """
        Proposta enviada não muda de conteúdo porque a empresa mudou de
        razão social depois. É o mesmo motivo do 'cargo' em uso_eventos.
        """
        h = usuario_adm["headers"]
        conta, opp = await cenario(client, h)
        proposta = (await client.post(f"/crm/oportunidades/{opp['id']}/propostas",
                                      json=corpo_proposta(), headers=h)).json()

        await client.patch(f"/crm/contas/{conta['id']}",
                           json={"razao_social": "Alfa Metais S.A."}, headers=h)

        lista = (await client.get(f"/crm/oportunidades/{opp['id']}/propostas",
                                  headers=h)).json()
        assert lista[0]["id"] == proposta["id"]
        assert lista[0]["cliente_razao_social"] == "Metalurgica Alfa LTDA"

    async def test_telefone_do_executivo_entra_do_cadastro(
        self, db_conn, client, usuario_adm
    ):
        h = usuario_adm["headers"]
        await client.put("/auth/perfil", json={"telefone": "+55 (11) 9 9571-3682"},
                         headers=h)
        _, opp = await cenario(client, h)
        body = (await client.post(f"/crm/oportunidades/{opp['id']}/propostas",
                                  json=corpo_proposta(), headers=h)).json()
        assert body["executivo_telefone"] == "+55 (11) 9 9571-3682"

    async def test_pode_assinar_por_outro_executivo(self, db_conn, client, usuario_adm):
        """
        O ADM monta a proposta enquanto o vendedor está em visita — e o
        slide precisa trazer o contato de quem vai atender a ligação.
        """
        h = usuario_adm["headers"]
        u = await criar_usuario(db_conn, client, "EV", "ev-proposta@teste.com")
        uid = str(await db_conn.fetchval(
            "SELECT id FROM usuarios WHERE email = $1", u["email"]))
        _, opp = await cenario(client, h)
        body = (await client.post(
            f"/crm/oportunidades/{opp['id']}/propostas",
            json=corpo_proposta(executivo_id=uid), headers=h,
        )).json()
        assert body["executivo_email"] == "ev-proposta@teste.com"
        assert body["criado_por_nome"] == "Test ADM"

    async def test_executivo_inexistente_e_422(self, db_conn, client, usuario_adm):
        _, opp = await cenario(client, usuario_adm["headers"])
        resp = await client.post(
            f"/crm/oportunidades/{opp['id']}/propostas",
            json=corpo_proposta(executivo_id=str(uuid.uuid4())),
            headers=usuario_adm["headers"],
        )
        assert resp.status_code == 422


# ── Validação pela API ───────────────────────────────────────────────

class TestValidacao:
    @pytest.mark.parametrize("troca", [
        {"vidas": 0},
        {"valor_por_vida": "0"},
        {"treinamentos": "-1"},
        {"escopo": []},
    ])
    async def test_recusa_numero_impossivel(self, db_conn, client, usuario_adm, troca):
        _, opp = await cenario(client, usuario_adm["headers"])
        resp = await client.post(
            f"/crm/oportunidades/{opp['id']}/propostas",
            json=corpo_proposta(**troca), headers=usuario_adm["headers"],
        )
        assert resp.status_code == 422

    async def test_validade_anterior_a_data(self, db_conn, client, usuario_adm):
        _, opp = await cenario(client, usuario_adm["headers"])
        resp = await client.post(
            f"/crm/oportunidades/{opp['id']}/propostas",
            json=corpo_proposta(validade="2026-08-25"),
            headers=usuario_adm["headers"],
        )
        assert resp.status_code == 422
        assert "anterior" in resp.json()["detail"]

    async def test_escopo_so_com_espacos(self, db_conn, client, usuario_adm):
        _, opp = await cenario(client, usuario_adm["headers"])
        resp = await client.post(
            f"/crm/oportunidades/{opp['id']}/propostas",
            json=corpo_proposta(escopo=["  ", ""]), headers=usuario_adm["headers"],
        )
        assert resp.status_code == 422


# ── Listagem e padrões ───────────────────────────────────────────────

class TestListagem:
    async def test_ordena_da_mais_nova_para_a_mais_antiga(
        self, db_conn, client, usuario_adm
    ):
        h = usuario_adm["headers"]
        _, opp = await cenario(client, h)
        for _ in range(3):
            await client.post(f"/crm/oportunidades/{opp['id']}/propostas",
                              json=corpo_proposta(), headers=h)
        lista = (await client.get(f"/crm/oportunidades/{opp['id']}/propostas",
                                  headers=h)).json()
        assert [p["versao"] for p in lista] == [3, 2, 1]

    async def test_oportunidade_sem_proposta_devolve_lista_vazia(
        self, db_conn, client, usuario_adm
    ):
        _, opp = await cenario(client, usuario_adm["headers"])
        resp = await client.get(f"/crm/oportunidades/{opp['id']}/propostas",
                                headers=usuario_adm["headers"])
        assert resp.status_code == 200
        assert resp.json() == []


class TestPadrao:
    async def test_traz_escopo_padrao_e_dados_do_logado(
        self, db_conn, client, usuario_adm
    ):
        _, opp = await cenario(client, usuario_adm["headers"])
        body = (await client.get(
            f"/crm/oportunidades/{opp['id']}/proposta-padrao",
            headers=usuario_adm["headers"],
        )).json()
        assert body["escopo_padrao"] == regras.ESCOPO_PADRAO
        assert body["executivo_nome"] == "Test ADM"
        assert body["cliente_razao_social"] == "Metalurgica Alfa LTDA"
        assert body["dias_validade"] == regras.DIAS_VALIDADE_PADRAO
        # Capacidades do servidor: a tela usa para avisar ANTES de o usuário
        # preencher o formulário inteiro.
        assert body["geracao_disponivel"] is True
        assert isinstance(body["pdf_disponivel"], bool)

    async def test_repete_a_ultima_proposta(self, db_conn, client, usuario_adm):
        """
        'Ajustar o desconto' é o caso comum: muda um valor, o resto continua
        igual. Formulário em branco obrigaria a redigitar tudo.
        """
        h = usuario_adm["headers"]
        _, opp = await cenario(client, h)
        await client.post(f"/crm/oportunidades/{opp['id']}/propostas",
                          json=corpo_proposta(vidas=77), headers=h)
        body = (await client.get(f"/crm/oportunidades/{opp['id']}/proposta-padrao",
                                 headers=h)).json()
        assert body["vidas"] == 77

    async def test_sem_proposta_anterior_vem_vazio(self, db_conn, client, usuario_adm):
        _, opp = await cenario(client, usuario_adm["headers"])
        body = (await client.get(f"/crm/oportunidades/{opp['id']}/proposta-padrao",
                                 headers=usuario_adm["headers"])).json()
        assert body["vidas"] is None


# ── Download ─────────────────────────────────────────────────────────

class TestDownload:
    async def test_pptx_sai_preenchido(self, db_conn, client, usuario_adm):
        """
        Monta o arquivo de verdade e confere que nenhum marcador sobrou.
        É o teste que pega modelo trocado sem as marcações.
        """
        from io import BytesIO
        import re
        from pptx import Presentation

        h = usuario_adm["headers"]
        _, opp = await cenario(client, h)
        proposta = (await client.post(f"/crm/oportunidades/{opp['id']}/propostas",
                                      json=corpo_proposta(), headers=h)).json()

        resp = await client.get(f"/crm/propostas/{proposta['id']}/arquivo",
                                params={"formato": "pptx"}, headers=h)
        assert resp.status_code == 200, resp.text
        assert "presentationml" in resp.headers["content-type"]
        assert opp["numero"] in resp.headers["content-disposition"]

        prs = Presentation(BytesIO(resp.content))
        texto = "\n".join(
            sh.text_frame.text
            for slide in prs.slides for sh in slide.shapes
            if sh.has_text_frame
        )
        assert not re.findall(r"\{\{\w+\}\}", texto), "sobrou marcador no arquivo"
        assert "Metalurgica Alfa LTDA" in texto
        assert "R$ 4.000,00" in texto
        assert "QTDE. VIDAS: 50" in texto
        assert "PGR - (NR-01)" in texto
        assert "05/09/2026" in texto

    async def test_pdf_sem_libreoffice_explica_o_que_fazer(
        self, db_conn, client, usuario_adm, monkeypatch
    ):
        """
        Falta de LibreOffice é problema de ambiente, não pedido inválido —
        503, e a mensagem precisa dizer que o PPTX continua funcionando.
        """
        h = usuario_adm["headers"]
        _, opp = await cenario(client, h)
        proposta = (await client.post(f"/crm/oportunidades/{opp['id']}/propostas",
                                      json=corpo_proposta(), headers=h)).json()

        monkeypatch.setattr(render, "libreoffice_disponivel", lambda: None)
        resp = await client.get(f"/crm/propostas/{proposta['id']}/arquivo",
                                params={"formato": "pdf"}, headers=h)
        assert resp.status_code == 503
        assert "PPTX" in resp.json()["detail"]

    async def test_sem_python_pptx_explica_o_que_instalar(
        self, db_conn, client, usuario_adm, monkeypatch
    ):
        """
        A biblioteca é importada tarde de propósito: faltar derruba só a
        proposta, não a API inteira. O 503 tem que dizer o comando.
        """
        h = usuario_adm["headers"]
        _, opp = await cenario(client, h)
        proposta = (await client.post(f"/crm/oportunidades/{opp['id']}/propostas",
                                      json=corpo_proposta(), headers=h)).json()

        def sem_biblioteca():
            raise render.BibliotecaIndisponivel(
                "python-pptx não está instalado no servidor. "
                "Instale com: pip install python-pptx==1.0.2"
            )
        monkeypatch.setattr(render, "_presentation", sem_biblioteca)

        resp = await client.get(f"/crm/propostas/{proposta['id']}/arquivo", headers=h)
        assert resp.status_code == 503
        assert "python-pptx" in resp.json()["detail"]

    async def test_formato_desconhecido_e_422(self, db_conn, client, usuario_adm):
        h = usuario_adm["headers"]
        _, opp = await cenario(client, h)
        proposta = (await client.post(f"/crm/oportunidades/{opp['id']}/propostas",
                                      json=corpo_proposta(), headers=h)).json()
        resp = await client.get(f"/crm/propostas/{proposta['id']}/arquivo",
                                params={"formato": "docx"}, headers=h)
        assert resp.status_code == 422

    async def test_proposta_inexistente(self, db_conn, client, usuario_adm):
        resp = await client.get(f"/crm/propostas/{uuid.uuid4()}/arquivo",
                                headers=usuario_adm["headers"])
        assert resp.status_code == 404


# ── Permissões ───────────────────────────────────────────────────────

class TestPermissoes:
    @pytest.mark.parametrize("cargo", ["Franqueado", "ADM", "EC", "SDR", "EV", "EP"])
    async def test_todo_cargo_valido_acessa(self, db_conn, client, cargo):
        """
        Proposta é operação de CRM: quem vende, gera. A tela é que decide
        quem vê o botão.
        """
        u = await criar_usuario(db_conn, client, cargo, f"{cargo.lower()}-prop@teste.com")
        conta = await nova_conta(client, u["headers"])
        opp = await nova_oportunidade(client, u["headers"], conta["id"])
        resp = await client.get(f"/crm/oportunidades/{opp['id']}/propostas",
                                headers=u["headers"])
        assert resp.status_code == 200

    async def test_cargo_extinto_403(self, db_conn, client, usuario_adm):
        u = await criar_usuario(db_conn, client, "Gerente", "gerente-prop@teste.com")
        _, opp = await cenario(client, usuario_adm["headers"])
        resp = await client.get(f"/crm/oportunidades/{opp['id']}/propostas",
                                headers=u["headers"])
        assert resp.status_code == 403

    async def test_sem_token_401(self, db_conn, client, usuario_adm):
        _, opp = await cenario(client, usuario_adm["headers"])
        assert (await client.get(
            f"/crm/oportunidades/{opp['id']}/propostas")).status_code == 401


# ── Telefone no perfil ───────────────────────────────────────────────

class TestPerfilTelefone:
    async def test_salva_e_volta_no_me(self, db_conn, client, usuario_adm):
        h = usuario_adm["headers"]
        assert (await client.get("/auth/me", headers=h)).json()["telefone"] is None

        resp = await client.put("/auth/perfil", json={"telefone": "11 90000-0000"},
                                headers=h)
        assert resp.status_code == 200
        assert (await client.get("/auth/me", headers=h)).json()["telefone"] == "11 90000-0000"

    async def test_string_vazia_limpa_o_campo(self, db_conn, client, usuario_adm):
        h = usuario_adm["headers"]
        await client.put("/auth/perfil", json={"telefone": "11 90000-0000"}, headers=h)
        await client.put("/auth/perfil", json={"telefone": "   "}, headers=h)
        assert (await client.get("/auth/me", headers=h)).json()["telefone"] is None

    async def test_sem_token_401(self, db_conn, client):
        assert (await client.put("/auth/perfil", json={"telefone": "x"})).status_code == 401
