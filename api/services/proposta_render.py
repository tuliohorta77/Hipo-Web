"""
HIPO — Preenchimento do modelo .pptx da proposta e conversão para PDF.

Separado de services/proposta.py de propósito: aqui há I/O (abre arquivo,
chama subprocesso) e dependência externa (python-pptx, LibreOffice). As
regras e a formatação ficam lá, testáveis sem nada disso.

## Como o modelo funciona

`api/templates/proposta_modelo.pptx` é o material da Controller MedSeg com
os campos variáveis trocados por marcadores `{{ASSIM}}`. Os slides 1 a 4
são institucionais e não têm marcador nenhum — o código nem os visita.

Cada marcador vive num run ÚNICO dentro do parágrafo. Isso não é detalhe:
o PowerPoint quebra texto em runs por corretor ortográfico e formatação, e
um `{{VALOR_VIDA}}` digitado à mão costuma virar três runs ('{{VALOR', '_',
'VIDA}}'), que nenhum replace de string encontra. O modelo foi gerado por
script justamente para garantir um run por marcador.

Se você trocar o modelo, mantenha os marcadores intactos e prefira colar
cada um de uma vez, sem editar letra por letra dentro dele.

## O escopo

`{{ESCOPO_ITEM}}` é o parágrafo-molde da lista. O código o clona uma vez
por item da proposta, preservando bullet, fonte e recuo, e remove o molde
no fim. Escrever os itens com '\\n' num run só perderia a marcação de lista.
"""
from __future__ import annotations

import copy
import os
import shutil
import subprocess
import tempfile
from io import BytesIO
from pathlib import Path

# python-pptx é importado DENTRO das funções, não aqui.
#
# Este módulo é alcançado pela cadeia main -> routers -> services, então um
# import no topo faz a API INTEIRA morrer se a biblioteca faltar: o serviço
# não sobe, e o HIPO fica fora do ar por causa de uma feature. Já aconteceu
# na primeira tentativa de deploy da 009 -- a suíte inteira de testes
# abortou no conftest, sem nenhum teste ter rodado.
#
# O deploy do CI faz rsync e reinicia; NÃO roda pip install. Com o import
# tardio, faltar a lib derruba só a proposta, com 503 e mensagem dizendo o
# que instalar.

# api/services/proposta_render.py -> api/templates/proposta_modelo.pptx
CAMINHO_MODELO = Path(__file__).resolve().parent.parent / "templates" / "proposta_modelo.pptx"

MARCADOR_ESCOPO = "{{ESCOPO_ITEM}}"

# Nomes do binário no PATH, conforme a distribuição.
BINARIOS_LIBREOFFICE = ("soffice", "libreoffice")

# A Amazon Linux 2023 NÃO tem LibreOffice nos repositórios (o
# libreoffice-impress do AL2 sumiu). Quando ele é instalado pelo tarball
# oficial da Document Foundation, o binário fica fora do PATH, em
# /opt/libreoffice<versão>/program/soffice — daí a busca por padrão de
# caminho além do `which`.
PADROES_LIBREOFFICE = (
    "/opt/libreoffice*/program/soffice",
    "/usr/lib64/libreoffice/program/soffice",
    "/opt/libreoffice*/program/soffice.bin",
)

# Conversão de 6 slides com imagens pesadas leva ~5s numa t3.medium fria.
# 120s é folga para o primeiro uso, quando o LibreOffice ainda monta o
# perfil do usuário.
TIMEOUT_PDF_S = 120


class BibliotecaIndisponivel(RuntimeError):
    """python-pptx não instalado no ambiente que está servindo."""


class ModeloIndisponivel(RuntimeError):
    """O .pptx do modelo não está onde deveria."""


class PdfIndisponivel(RuntimeError):
    """LibreOffice ausente ou falhou. A mensagem diz o que fazer no servidor."""


# ── PPTX ─────────────────────────────────────────────────────────────

def _presentation():
    """Importa python-pptx na hora do uso. Ver a nota no topo do módulo."""
    try:
        from pptx import Presentation
    except ImportError as erro:
        raise BibliotecaIndisponivel(
            "python-pptx não está instalado no servidor, então a proposta não "
            "pode ser gerada. Instale com: "
            "sudo -iu hipo python3 -m pip install --user python-pptx==1.0.2"
        ) from erro
    return Presentation


def pptx_disponivel() -> bool:
    """Usado pela tela para avisar antes de o usuário preencher o formulário."""
    try:
        import pptx  # noqa: F401
    except ImportError:
        return False
    return True


def _substituir_no_texto(shape, mapa: dict[str, str]) -> None:
    """
    Troca marcadores run a run, preservando a formatação de cada um.

    Percorre runs em vez do texto do shape inteiro porque atribuir
    `text_frame.text` apaga toda a formatação do quadro — o slide sairia
    com a fonte padrão do tema no lugar da tipografia da marca.
    """
    if not shape.has_text_frame:
        return
    for par in shape.text_frame.paragraphs:
        for run in par.runs:
            texto = run.text
            if "{{" not in texto:
                continue
            for marcador, valor in mapa.items():
                if marcador in texto:
                    texto = texto.replace(marcador, valor)
            run.text = texto


def _preencher_escopo(slide, itens: list[str]) -> bool:
    """
    Clona o parágrafo-molde uma vez por item. Devolve True se achou o molde.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        molde = None
        for par in shape.text_frame.paragraphs:
            if MARCADOR_ESCOPO in par.text:
                molde = par
                break
        if molde is None:
            continue

        pai = molde._p.getparent()
        for item in itens:
            novo = copy.deepcopy(molde._p)
            pai.insert(list(pai).index(molde._p), novo)
            # Reencontra o parágrafo recém-inserido pela árvore XML: o
            # objeto Paragraph do python-pptx é um wrapper, e o índice na
            # lista de paragraphs muda a cada inserção.
            from pptx.text.text import _Paragraph  # import local: detalhe interno
            _Paragraph(novo, shape.text_frame).runs[0].text = item

        pai.remove(molde._p)
        return True
    return False


def montar_pptx(
    substituicoes: dict[str, str],
    escopo: list[str],
    caminho_modelo: Path | str | None = None,
) -> bytes:
    """
    Devolve o .pptx preenchido, em memória.

    Não grava em disco: o arquivo é servido direto na resposta HTTP, e
    escrever num diretório temporário só criaria lixo para limpar (e uma
    corrida entre dois vendedores gerando ao mesmo tempo).
    """
    modelo = Path(caminho_modelo or CAMINHO_MODELO)
    if not modelo.is_file():
        raise ModeloIndisponivel(
            f"Modelo da proposta não encontrado em {modelo}. "
            "Ele é versionado em api/templates/ — confira se o deploy copiou a pasta."
        )

    Presentation = _presentation()
    prs = Presentation(str(modelo))

    for slide in prs.slides:
        _preencher_escopo(slide, escopo)
        for shape in slide.shapes:
            _substituir_no_texto(shape, substituicoes)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ── PDF ──────────────────────────────────────────────────────────────

def libreoffice_disponivel() -> str | None:
    """Caminho do binário, ou None. Usado pela tela para esconder o botão."""
    for nome in BINARIOS_LIBREOFFICE:
        caminho = shutil.which(nome)
        if caminho:
            return caminho

    # Instalação por tarball não põe nada no PATH.
    import glob
    for padrao in PADROES_LIBREOFFICE:
        for achado in sorted(glob.glob(padrao), reverse=True):
            if os.access(achado, os.X_OK):
                return achado
    return None


def para_pdf(pptx: bytes) -> bytes:
    """
    Converte via LibreOffice headless.

    Roda num diretório temporário próprio, com HOME apontando para ele: sem
    isso o soffice tenta escrever o perfil em /home/hipo e, se dois pedidos
    chegam juntos, o segundo morre disputando o mesmo lock — falha
    intermitente que só aparece quando dois vendedores geram ao mesmo tempo.
    """
    binario = libreoffice_disponivel()
    if not binario:
        raise PdfIndisponivel(
            "LibreOffice não está instalado no servidor, então o PDF não pode "
            "ser gerado. O PPTX continua funcionando normalmente — baixe e "
            "exporte pelo PowerPoint. Para habilitar o PDF, rode o "
            "deploy-009 com -InstalarLibreOffice (a Amazon Linux 2023 não "
            "tem o pacote nos repositórios; a instalação é pelo tarball "
            "oficial)."
        )

    with tempfile.TemporaryDirectory(prefix="hipo-proposta-") as tmp:
        entrada = os.path.join(tmp, "proposta.pptx")
        with open(entrada, "wb") as f:
            f.write(pptx)

        ambiente = dict(os.environ, HOME=tmp)
        try:
            resultado = subprocess.run(
                [binario, "--headless", "--norestore", "--invisible",
                 "--convert-to", "pdf", "--outdir", tmp, entrada],
                capture_output=True, timeout=TIMEOUT_PDF_S, env=ambiente,
            )
        except subprocess.TimeoutExpired as exc:
            raise PdfIndisponivel(
                f"A conversão para PDF passou de {TIMEOUT_PDF_S}s e foi "
                "interrompida. Baixe o PPTX e converta no PowerPoint."
            ) from exc

        saida = os.path.join(tmp, "proposta.pdf")
        if not os.path.exists(saida):
            erro = (resultado.stderr or b"").decode("utf-8", "replace")[:300]
            raise PdfIndisponivel(
                "O LibreOffice não produziu o PDF. Baixe o PPTX e converta no "
                f"PowerPoint. Detalhe do servidor: {erro or 'sem mensagem'}"
            )
        with open(saida, "rb") as f:
            return f.read()
