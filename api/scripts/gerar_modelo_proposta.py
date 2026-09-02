"""
Gera api/templates/proposta_modelo.pptx a partir do material original da
Controller MedSeg.

Rode quando o marketing entregar uma arte nova:

    python -m scripts.gerar_modelo_proposta "PROPOSTA COMERCIAL ... .pptx"

## Por que um script, e não editar o .pptx à mão

Cada marcador precisa viver num RUN ÚNICO. O PowerPoint quebra texto em
runs por corretor ortográfico e formatação — digitar `{{VALOR_VIDA}}` na
tela costuma produzir três runs ('{{VALOR', '_', 'VIDA}}'), e nenhum
replace de string encontra isso. O script consolida cada parágrafo em um
run só, herdando a formatação do primeiro, e escreve o marcador inteiro.

## Se a arte mudar de posição

Os shapes são localizados por ID, que é estável dentro de um arquivo mas
NÃO entre arquivos diferentes. Numa arte nova, rode o inspetor abaixo,
confira os IDs e ajuste os mapas:

    python -m scripts.gerar_modelo_proposta --inspecionar arquivo.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

RAIZ = Path(__file__).resolve().parent.parent
DESTINO = RAIZ / "templates" / "proposta_modelo.pptx"

# Slide 5 (índice 4): escopo e investimento.
SLIDE_ESCOPO = 4
# Slide 6 (índice 5): cliente, executivo, data e validade.
SLIDE_FECHAMENTO = 5

# shape_id -> texto com marcadores.
TEXTOS_SLIDE_5 = {
    16: "QTDE. VIDAS: {{VIDAS}}",
    14: "{{VALOR_VIDA}}",
    19: "{{MENSALIDADE}}",
    20: "{{TREINAMENTOS}}",
    22: "{{LAUDOS}}",
    24: "{{INVESTIMENTO}}",
    15: (
        "Assessoria permanente em segurança e Medicina do Trabalho enquanto "
        "vigorar o contrato*** Os exames complementares (quando necessários) "
        "terão seus valores acertados de acordo com a tabela vigente na data "
        "de realização. A mensalidade apresentada considera o número de vidas "
        "informado nesta proposta. Para cada vida adicional, será acrescido o "
        "valor de {{VALOR_VIDA}} mensais."
    ),
}

TEXTOS_SLIDE_6 = {
    25: "{{EXECUTIVO_NOME}}",
    27: "{{EXECUTIVO_EMAIL}}",
    29: "{{EXECUTIVO_TELEFONE}}",
    30: "{{CLIENTE}}",
    31: "{{CIDADE}}, {{DATA_EXTENSO}}",
    32: "Proposta válida até {{VALIDADE}}",
}

# A caixa de texto do escopo, cujo primeiro parágrafo vira o molde da lista.
SHAPE_ESCOPO = 8

# RÓTULOS do quadro de investimento — "Valor por vida ....", "Mensalidade
# ....", "INVESTIMENTO". As caixas foram dimensionadas no PowerPoint com a
# Codec Pro real; o LibreOffice mede um pixel a mais e quebra a linha,
# partindo "INVESTIMENTO" ao meio no PDF. Sem quebra, o texto transborda em
# vez de dobrar — e como ele já cabe na caixa, o resultado fica igual nos
# dois programas.
#
# Os VALORES (14, 19, 20, 22, 24) ficam de fora de propósito: são
# centralizados em caixas largas, e sem quebra cresceriam para os lados
# invadindo a coluna vizinha.
SEM_QUEBRA_SLIDE_5 = (17, 18, 21, 23, 25)


def _shape(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape {shape_id} não existe neste slide")


def _run_unico(paragrafo, texto: str) -> None:
    """Deixa o parágrafo com um run só, herdando a formatação do primeiro."""
    runs = paragrafo.runs
    if not runs:
        raise ValueError("parágrafo sem runs — nada de onde herdar formatação")
    runs[0].text = texto
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def inspecionar(caminho: str) -> None:
    """Lista shapes com texto, por slide, para conferir IDs numa arte nova."""
    prs = Presentation(caminho)
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n=== SLIDE {i} ===")
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                texto = sh.text_frame.text.strip().replace("\n", " / ")[:90]
                print(f"  [{sh.shape_id:>3}] {sh.name:<16} {texto}")


def gerar(origem: str, destino: Path = DESTINO) -> Path:
    prs = Presentation(origem)
    s5 = prs.slides[SLIDE_ESCOPO]
    s6 = prs.slides[SLIDE_FECHAMENTO]

    # Escopo: mantém só o primeiro parágrafo, que vira o molde clonado pelo
    # render uma vez por item.
    quadro = _shape(s5, SHAPE_ESCOPO).text_frame
    _run_unico(quadro.paragraphs[0], "{{ESCOPO_ITEM}}")
    for par in list(quadro.paragraphs[1:]):
        par._p.getparent().remove(par._p)

    for sid, texto in TEXTOS_SLIDE_5.items():
        _run_unico(_shape(s5, sid).text_frame.paragraphs[0], texto)
    for sid, texto in TEXTOS_SLIDE_6.items():
        _run_unico(_shape(s6, sid).text_frame.paragraphs[0], texto)

    for sid in SEM_QUEBRA_SLIDE_5:
        _shape(s5, sid).text_frame.word_wrap = False

    destino.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(destino))
    return destino


def main() -> int:
    args = sys.argv[1:]
    if not args:
        print(__doc__)
        return 1
    if args[0] == "--inspecionar":
        inspecionar(args[1])
        return 0

    destino = gerar(args[0])
    tamanho = destino.stat().st_size / 1024 / 1024
    print(f"gerado: {destino} ({tamanho:.1f} MB)")

    # Confere que todo marcador esperado sobreviveu. Modelo salvo sem um
    # marcador gera proposta com campo em branco, e ninguém percebe até o
    # cliente receber.
    import re
    achados = set()
    for slide in Presentation(str(destino)).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                achados.update(re.findall(r"\{\{\w+\}\}", sh.text_frame.text))
    esperados = {
        "{{ESCOPO_ITEM}}", "{{VIDAS}}", "{{VALOR_VIDA}}", "{{MENSALIDADE}}",
        "{{TREINAMENTOS}}", "{{LAUDOS}}", "{{INVESTIMENTO}}", "{{CLIENTE}}",
        "{{EXECUTIVO_NOME}}", "{{EXECUTIVO_EMAIL}}", "{{EXECUTIVO_TELEFONE}}",
        "{{CIDADE}}", "{{DATA_EXTENSO}}", "{{VALIDADE}}",
    }
    faltando = esperados - achados
    if faltando:
        print("ERRO: marcadores ausentes:", sorted(faltando))
        return 1
    print(f"{len(achados)} marcadores conferidos")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
